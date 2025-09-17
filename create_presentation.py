#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

def create_agile_point_presentation():
    # Criar nova apresentação
    prs = Presentation()
    
    # Definir cores da marca (mais vibrantes)
    primary_color = RGBColor(17, 17, 17)  # #111111
    secondary_color = RGBColor(100, 100, 100)  # #646464
    accent_color = RGBColor(0, 123, 255)  # Azul
    background_color = RGBColor(248, 250, 252)  # Fundo claro
    card_color = RGBColor(255, 255, 255)  # Branco para cards
    
    # SLIDE 1 - CAPA (Fundo escuro moderno)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    
    # Fundo escuro
    background = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = primary_color
    background.line.fill.background()
    
    # Retângulo branco para logo (como no HTML)
    logo_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1), Inches(2), Inches(0.5))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = card_color
    logo_box.line.fill.background()
    
    # Título principal
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "TRANSFORMAMOS DESAFIOS EM REALIDADE DIGITAL"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = card_color
    
    # Subtítulo
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Apresentação Corporativa 2025"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = secondary_color
    
    # SLIDE 2 - QUEM SOMOS (Fundo escuro como no HTML)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    
    # Fundo escuro
    background = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(40, 40, 40)  # Cinza escuro
    background.line.fill.background()
    
    # Título
    title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "QUEM SOMOS?"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.LEFT
    title_para.font.size = Pt(42)
    title_para.font.bold = True
    title_para.font.color.rgb = card_color
    
    # Conteúdo principal
    content_box = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(2.5))
    content_frame = content_box.text_frame
    content_frame.text = """Na Agile Point, transformamos seus desafios em realidade digital. Com mais de 17 anos de experiência, nossa missão é clara: simplificar o complexo.

Desenvolvemos soluções inteligentes e eficientes que otimizam seus processos e resultados que impulsionam seu crescimento."""
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = card_color
        paragraph.space_after = Pt(16)
    
    # Cards com métricas (como no HTML)
    # Card 1 - 17+ Anos
    card1 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.5), Inches(2.5), Inches(1.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(60, 60, 60)
    card1.line.fill.background()
    
    metric1_box = slide2.shapes.add_textbox(Inches(1.1), Inches(4.6), Inches(2.3), Inches(0.8))
    metric1_frame = metric1_box.text_frame
    metric1_frame.text = "17+"
    metric1_para = metric1_frame.paragraphs[0]
    metric1_para.alignment = PP_ALIGN.CENTER
    metric1_para.font.size = Pt(36)
    metric1_para.font.bold = True
    metric1_para.font.color.rgb = card_color
    
    desc1_box = slide2.shapes.add_textbox(Inches(1.1), Inches(5.4), Inches(2.3), Inches(0.5))
    desc1_frame = desc1_box.text_frame
    desc1_frame.text = "Anos de Experiência"
    desc1_para = desc1_frame.paragraphs[0]
    desc1_para.alignment = PP_ALIGN.CENTER
    desc1_para.font.size = Pt(14)
    desc1_para.font.color.rgb = card_color
    
    # Card 2 - 100+ Projetos
    card2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(4.5), Inches(2.5), Inches(1.5))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(60, 60, 60)
    card2.line.fill.background()
    
    metric2_box = slide2.shapes.add_textbox(Inches(3.9), Inches(4.6), Inches(2.3), Inches(0.8))
    metric2_frame = metric2_box.text_frame
    metric2_frame.text = "100+"
    metric2_para = metric2_frame.paragraphs[0]
    metric2_para.alignment = PP_ALIGN.CENTER
    metric2_para.font.size = Pt(36)
    metric2_para.font.bold = True
    metric2_para.font.color.rgb = card_color
    
    desc2_box = slide2.shapes.add_textbox(Inches(3.9), Inches(5.4), Inches(2.3), Inches(0.5))
    desc2_frame = desc2_box.text_frame
    desc2_frame.text = "Projetos Entregues"
    desc2_para = desc2_frame.paragraphs[0]
    desc2_para.alignment = PP_ALIGN.CENTER
    desc2_para.font.size = Pt(14)
    desc2_para.font.color.rgb = card_color
    
    # Card 3 - 100% Satisfação
    card3 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(4.5), Inches(2.5), Inches(1.5))
    card3.fill.solid()
    card3.fill.fore_color.rgb = RGBColor(60, 60, 60)
    card3.line.fill.background()
    
    metric3_box = slide2.shapes.add_textbox(Inches(6.7), Inches(4.6), Inches(2.3), Inches(0.8))
    metric3_frame = metric3_box.text_frame
    metric3_frame.text = "100%"
    metric3_para = metric3_frame.paragraphs[0]
    metric3_para.alignment = PP_ALIGN.CENTER
    metric3_para.font.size = Pt(36)
    metric3_para.font.bold = True
    metric3_para.font.color.rgb = card_color
    
    desc3_box = slide2.shapes.add_textbox(Inches(6.7), Inches(5.4), Inches(2.3), Inches(0.5))
    desc3_frame = desc3_box.text_frame
    desc3_frame.text = "Satisfação do Cliente"
    desc3_para = desc3_frame.paragraphs[0]
    desc3_para.alignment = PP_ALIGN.CENTER
    desc3_para.font.size = Pt(14)
    desc3_para.font.color.rgb = card_color
    
    # SLIDE 3 - NOSSOS SERVIÇOS
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide3.shapes.title.text = "NOSSOS SERVIÇOS"
    slide3.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide3.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide3.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide3.placeholders[1]
    content.text = """• Soluções Web e Mobile
  Criamos experiências digitais modernas e intuitivas que engajam seus clientes e expandem seu alcance no mercado.

• Automação de Processos (RPA)
  Automatizamos tarefas repetitivas, liberando sua equipe para focar em atividades estratégicas e reduzindo custos operacionais.

• Desenvolvimento Customizado
  Construímos soluções sob medida que se adaptam perfeitamente às suas necessidades, garantindo escalabilidade e eficiência.

• Consultoria em Tecnologia
  Orientação estratégica para transformação digital e otimização de processos."""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(8)
    
    # SLIDE 4 - CASES DE SUCESSO
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide4.shapes.title.text = "CASES DE SUCESSO"
    slide4.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide4.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide4.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide4.placeholders[1]
    content.text = """[Adicione aqui casos reais da sua empresa]

Exemplo:
• Empresa X - Sistema de Gestão
  Redução de 60% no tempo de processos
  ROI de 300% em 6 meses

• Empresa Y - Automação RPA
  Economia de 40 horas/semana
  Eliminação de erros manuais"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 5 - ESCOPO DO PROJETO
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide5.shapes.title.text = "ESCOPO DO PROJETO"
    slide5.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide5.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide5.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide5.placeholders[1]
    content.text = """[Defina aqui o escopo específico para cada cliente]

Exemplo:
• Análise de Requisitos
• Desenvolvimento da Solução
• Testes e Validação
• Implementação
• Treinamento da Equipe
• Suporte Pós-Implementação"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 6 - STACK TECNOLÓGICA
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide6.shapes.title.text = "STACK TECNOLÓGICA"
    slide6.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide6.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide6.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide6.placeholders[1]
    content.text = """Frontend:
• React.js / Next.js
• TypeScript
• Tailwind CSS

Backend:
• Python / Node.js
• APIs RESTful
• Banco de Dados

Ferramentas:
• Git / GitHub
• Docker
• CI/CD"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 7 - REQUISITOS TÉCNICOS
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide7.shapes.title.text = "REQUISITOS TÉCNICOS"
    slide7.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide7.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide7.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide7.placeholders[1]
    content.text = """Servidor:
• Mínimo 4GB RAM
• 50GB de armazenamento
• SSL Certificate

Domínio:
• Registro de domínio
• Configuração DNS

Backup:
• Backup automático diário
• Restore em caso de falhas"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 8 - ESTILIZAÇÃO E COMPONENTES
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide8.shapes.title.text = "ESTILIZAÇÃO E COMPONENTES"
    slide8.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide8.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide8.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide8.placeholders[1]
    content.text = """Design System:
• Paleta de cores personalizada
• Tipografia (Anton + Lato)
• Componentes reutilizáveis

Responsividade:
• Mobile First
• Tablet e Desktop
• Cross-browser compatibility

Acessibilidade:
• WCAG 2.1 AA
• Navegação por teclado
• Screen readers"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 9 - HOSPEDAGEM E BANCO DE DADOS
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide9.shapes.title.text = "HOSPEDAGEM E BANCO DE DADOS"
    slide9.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide9.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide9.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide9.placeholders[1]
    content.text = """Hospedagem:
• Vercel / AWS / Google Cloud
• CDN global
• SSL automático

Banco de Dados:
• PostgreSQL / MongoDB
• Backup automático
• Monitoramento 24/7

Segurança:
• Firewall configurado
• Criptografia de dados
• Logs de auditoria"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 10 - EQUIPE E TIMELINE
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide10.shapes.title.text = "EQUIPE E TIMELINE"
    slide10.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide10.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide10.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide10.placeholders[1]
    content.text = """Equipe:
• 1 Project Manager
• 2 Desenvolvedores
• 1 Designer UX/UI
• 1 QA Tester

Timeline:
• Semana 1-2: Análise e Planejamento
• Semana 3-6: Desenvolvimento
• Semana 7: Testes e Ajustes
• Semana 8: Deploy e Treinamento"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 11 - INVESTIMENTO E FORMA DE PAGAMENTO
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide11.shapes.title.text = "INVESTIMENTO E FORMA DE PAGAMENTO"
    slide11.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide11.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide11.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide11.placeholders[1]
    content.text = """Investimento:
• [Valor total do projeto]

Forma de Pagamento:
• 30% na assinatura do contrato
• 40% na entrega da primeira versão
• 30% na entrega final

Garantia:
• 90 dias de suporte gratuito
• Correções de bugs inclusas
• Treinamento da equipe"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 12 - CRONOGRAMA MACRO
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide12.shapes.title.text = "CRONOGRAMA MACRO"
    slide12.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide12.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide12.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide12.placeholders[1]
    content.text = """[Adicione aqui um cronograma visual com as principais etapas]

Exemplo:
Mês 1: Análise e Planejamento
Mês 2: Desenvolvimento
Mês 3: Testes e Ajustes
Mês 4: Deploy e Treinamento"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 13 - PRÓXIMOS PASSOS
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide13.shapes.title.text = "PRÓXIMOS PASSOS"
    slide13.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide13.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide13.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide13.placeholders[1]
    content.text = """1. Aprovação da proposta
2. Assinatura do contrato
3. Kickoff do projeto
4. Início do desenvolvimento

Contato:
contatoagilepoint@gmail.com
[Seu telefone]"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = secondary_color
        paragraph.space_after = Pt(12)
    
    # SLIDE 14 - OBRIGADO
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    
    # Título
    title_box = slide14.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "OBRIGADO!"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Subtítulo
    subtitle_box = slide14.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Vamos transformar seus desafios em realidade digital?"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = secondary_color
    
    # Contato
    contact_box = slide14.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "contatoagilepoint@gmail.com"
    contact_para = contact_frame.paragraphs[0]
    contact_para.alignment = PP_ALIGN.CENTER
    contact_para.font.size = Pt(18)
    contact_para.font.color.rgb = accent_color
    
    # Salvar apresentação
    prs.save('AGILE_POINT_PRESENTATION.pptx')
    print("✅ Apresentação criada com sucesso: AGILE_POINT_PRESENTATION.pptx")

if __name__ == "__main__":
    create_agile_point_presentation()